"""Markdown → PPTX 演示文稿生成服务（V1.2.1）。

将 PPT 制作功能生成的最终 Markdown 内容按「幻灯片 Markdown 约定」渲染为 16:9 的 pptx：

- 封面页：以传入的 title 为主标题（居中大字 + 顶部强调色条）。
- 内容页：`##`/`#`（heading1）→ 新建一页并作为页标题；`###`+（heading2）→ 页内小标题；
  `- ` / `1. ` / 段落 / 引用 / 代码 → 当前页要点文本框；表格 → add_table；图片 → add_picture；
  独占一行 `---` → 强制换页。
- 中文字体：所有 run 同时设置西文（a:latin）与中文（a:ea）字体，避免中文显示为默认西文字体。
- 溢出保护：单页要点过多/过长时自动续页（标题加「（续）」前缀）。

复用 docx_service._parse_blocks 解析 Markdown（同一套块模型）。
"""
import os
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .docx_service import _parse_blocks, _resolve_image_path

logger = logging.getLogger("app.pptx")

# ---------------- 版式与颜色 ----------------
SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

ACCENT = RGBColor(0x1F, 0x5C, 0x99)      # 深蓝强调色
TITLE_COLOR = RGBColor(0x1F, 0x2D, 0x3D)  # 页标题深色
BODY_COLOR = RGBColor(0x33, 0x3A, 0x45)   # 正文深灰
MUTED_COLOR = RGBColor(0x6B, 0x72, 0x80)  # 次要文字

FONT_CN = "微软雅黑"  # Microsoft YaHei

BODY_PT = 18
SUBHEAD_PT = 20
TITLE_PT = 28
COVER_PT = 40
CODE_PT = 14
TABLE_PT = 12

# 单页内容溢出保护阈值
MAX_PARAS = 9
MAX_CHARS = 900


def _set_run_font(run, size_pt: float, bold: bool = False, color=RGBColor(0, 0, 0),
                  italic: bool = False, font_cn: str = FONT_CN) -> None:
    """设置 run 字号/加粗/颜色，并同时指定西文与中文（eastAsia）字体。"""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_cn  # 设置 <a:latin typeface>
    rPr = run._r.get_or_add_rPr()
    # 追加 <a:ea typeface> 指定中文字体（schema 顺序：latin 之后）
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {"typeface": font_cn})
        latin = rPr.find(qn("a:latin"))
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    else:
        ea.set("typeface", font_cn)


def _add_textbox(slide, left, top, width, height, word_wrap: bool = True):
    """新建文本框，返回 text_frame。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    return tf


def _add_cover(prs, blank_layout, title_text: str) -> None:
    """封面页：顶部强调色条 + 居中大标题 + 副标题。"""
    slide = prs.slides.add_slide(blank_layout)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    tf = _add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.6))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text.strip() or "演示文稿"
    _set_run_font(run, COVER_PT, bold=True, color=ACCENT)

    stf = _add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.8))
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    srun = sp.add_run()
    srun.text = "规范智能问答助手 · 自动生成"
    _set_run_font(srun, 18, color=MUTED_COLOR)


def _add_content_slide(prs, blank_layout, slide_title: str) -> None:
    """新增内容页：顶部页标题 + 强调色线 + 正文文本框。返回 (slide, text_frame)。"""
    slide = prs.slides.add_slide(blank_layout)

    ttf = _add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9))
    tp = ttf.paragraphs[0]
    trun = tp.add_run()
    trun.text = slide_title
    _set_run_font(trun, TITLE_PT, bold=True, color=TITLE_COLOR)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(12.3), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    ctf = _add_textbox(slide, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.4))
    return slide, ctf


class _Deck:
    """渲染状态：当前页与正文文本框。"""

    def __init__(self, prs, blank_layout):
        self.prs = prs
        self.blank = blank_layout
        self.tf = None
        self.title = ""
        self.paras = 0
        self.chars = 0

    def new_slide(self, slide_title: str):
        _, self.tf = _add_content_slide(self.prs, self.blank, slide_title)
        self.title = slide_title
        self.paras = 0
        self.chars = 0

    def ensure(self):
        """确保存在当前内容页（首个内容块之前可能还没有页）。"""
        if self.tf is None:
            self.new_slide("")

    def _maybe_spill(self):
        """单页内容过多时续页。"""
        if self.tf is not None and (self.paras >= MAX_PARAS or self.chars >= MAX_CHARS):
            self.new_slide(f"（续）{self.title}" if self.title else "（续）")

    def add_paragraph(self, text: str, prefix: str = "", size: float = BODY_PT,
                      bold: bool = False, color=BODY_COLOR) -> None:
        """往当前页追加一段（自动续页保护）。"""
        self.ensure()
        self._maybe_spill()
        tf = self.tf
        p = tf.paragraphs[0] if self.paras == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"{prefix}{text}"
        _set_run_font(run, size, bold=bold, color=color)
        self.paras += 1
        self.chars += len(run.text)

    def add_table(self, rows: list[list[str]]) -> None:
        """往当前页底部添加 Markdown 表格。"""
        self.ensure()
        slide = self.prs.slides[-1]
        nrows, ncols = len(rows), max(len(r) for r in rows) if rows else 0
        if nrows == 0 or ncols == 0:
            return
        left, top, width, height = Inches(1.0), Inches(1.9), Inches(11.3), Inches(0.4 * nrows)
        graphic = slide.shapes.add_table(nrows, ncols, left, top, width, height)
        table = graphic.table
        for ri, row in enumerate(rows):
            for ci in range(ncols):
                cell = table.cell(ri, ci)
                cell.text = ""
                para = cell.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = row[ci].strip() if ci < len(row) else ""
                _set_run_font(run, TABLE_PT, bold=(ri == 0))
        self.paras += nrows
        self.chars += 200  # 表格粗略计字，触发续页阈值

    def add_image(self, url: str, alt: str) -> None:
        """往当前页插入图片（居中，宽度 6 英寸）；缺失时输出文字占位。"""
        self.ensure()
        slide = self.prs.slides[-1]
        abs_path = _resolve_image_path(url)
        if abs_path and os.path.isfile(abs_path):
            try:
                pic_w = Inches(6)
                left = (SLIDE_W - pic_w) // 2
                slide.shapes.add_picture(abs_path, left, Inches(1.9), width=pic_w)
                self.paras += 3
                return
            except Exception as e:  # noqa: BLE001
                logger.warning("插入图片失败 %s: %s", url, e)
        self.add_paragraph(f"[图片：{alt or url}]", size=BODY_PT, color=MUTED_COLOR)


def generate_pptx(content_markdown: str, title: str, out_path: str) -> str:
    """生成 16:9 的 pptx 演示文稿，返回文件路径。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    _add_cover(prs, blank, title)

    deck = _Deck(prs, blank)
    for block in _parse_blocks(content_markdown):
        kind = block["type"]
        if kind == "heading1":
            # `#`/`##` → 新一页
            deck.new_slide(block["text"])
        elif kind == "heading2":
            # `###`+ → 当前页内小标题
            deck.add_paragraph(block["text"], size=SUBHEAD_PT, bold=True, color=ACCENT)
        elif kind == "paragraph":
            text = (block["text"] or "").strip()
            if text in ("---", "***", "___"):
                # 独立分隔行 → 强制换页（下一块内容自动开新页）
                deck.tf = None
            elif text:
                deck.add_paragraph(text)
        elif kind == "bullet":
            deck.add_paragraph(block["text"], prefix="• ")
        elif kind == "ordered":
            deck.add_paragraph(block["text"], prefix=f"{block['num']}. ")
        elif kind == "quote":
            deck.add_paragraph(block["text"], color=MUTED_COLOR, bold=False)
        elif kind == "code":
            deck.add_paragraph(block["text"], size=CODE_PT)
        elif kind == "table":
            deck.add_table(block["rows"])
        elif kind == "image":
            deck.add_image(block["url"], block["alt"])

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    logger.info("pptx 演示文稿已生成: %s", out_path)
    return out_path
